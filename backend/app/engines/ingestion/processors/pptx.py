import io
from pptx import Presentation
from ..base import BaseSourceProcessor, ExtractionResult
from ..validators import validate_file


class PPTXProcessor(BaseSourceProcessor):
    source_type = "pptx"

    def validate(self, filename: str = "", content: bytes = b"", **kwargs) -> None:
        validate_file(filename, content)

    def extract(self, filename: str = "", content: bytes = b"", filepath: str = "", **kwargs) -> ExtractionResult:
        if filepath:
            prs = Presentation(filepath)
        else:
            prs = Presentation(io.BytesIO(content))

        slides = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- Slide {idx} ---"]
            if slide.shapes.title and slide.shapes.title.text:
                slide_parts.append(f"# {slide.shapes.title.text}")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        slide_parts.append(shape.text.strip())

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"[Speaker Notes] {notes}")

            slides.append("\n".join(slide_parts))

        full_text = "\n\n".join(slides)
        metadata = {"slide_count": len(prs.slides)}
        title = filename or "PowerPoint Presentation"
        return ExtractionResult(text=full_text, metadata=metadata, title=title)
